import os
import io
import json
import time
import logging
import requests
from pathlib import Path
from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Optional
import fitz
from docx import Document
from pptx import Presentation
from docx.shared import Pt

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S"
)
log = logging.getLogger(__name__)

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── Stats ────────────────────────────────────────────────────────────────────
STATS_FILE = Path(__file__).parent / "stats.json"

def load_stats():
    if STATS_FILE.exists():
        try:
            return json.loads(STATS_FILE.read_text())
        except Exception:
            pass
    return {"visits": 0, "translations": 0}

def save_stats(stats):
    STATS_FILE.write_text(json.dumps(stats))

def increment(key):
    stats = load_stats()
    stats[key] = stats.get(key, 0) + 1
    save_stats(stats)
    return stats

# ── AI config ────────────────────────────────────────────────────────────────
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
GROQ_API_KEY   = os.environ.get("Groq_API_KEY")   # именно так как в Railway
OLLAMA_URL     = os.environ.get("OLLAMA_URL", "http://localhost:11434/api/generate")
OLLAMA_MODEL   = os.environ.get("OLLAMA_MODEL", "qwen3.5:4b-mlx")

PROMPT_TEMPLATE = """Сен — қазақ тілінің кәсіби аудармашысысың. Мемлекеттік құжаттарды аударуда тәжірибелісің.

Аудару талаптары:
- Тек кириллица қарпін қолдан
- Ресми мемлекеттік стильді сақта
- Орыс сөздерінің орнына қазақ баламаларын қолдан:
  * «критерии» → «өлшемшарттар»
  * «аттестация нәтижелері» → «аттестаттау қорытындылары»
  * «продвижение» → «ілгерілету»
- Заңдық терминдерді дәл аудар
- Мәтін құрылымы мен нөмірлеуін сақта
- Тек аударманы жаз, түсініктеме берме

Аударылатын мәтін:
{text}"""

def translate_to_kazakh(text: str, model: str = "gemini") -> str:
    prompt = PROMPT_TEMPLATE.format(text=text)
    chars = len(text)
    t0 = time.time()

    log.info(f"[TRANSLATE] модель={model} символов={chars}")

    if model == "groq":
        if not GROQ_API_KEY:
            raise HTTPException(503, "GROQ_API_KEY не настроен на сервере")
        log.info("[GROQ] отправка запроса...")
        resp = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"},
            json={
                "model": "llama-3.3-70b-versatile",
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.3,
            },
            timeout=120
        )
        resp.raise_for_status()
        result = resp.json()["choices"][0]["message"]["content"]

    elif model == "ollama":
        log.info(f"[OLLAMA] отправка на {OLLAMA_URL} модель={OLLAMA_MODEL}")
        resp = requests.post(OLLAMA_URL, json={
            "model": OLLAMA_MODEL,
            "prompt": prompt,
            "stream": False
        }, headers={"ngrok-skip-browser-warning": "true"}, timeout=300)
        resp.raise_for_status()
        result = resp.json()["response"]

    else:  # gemini (default)
        if not GEMINI_API_KEY:
            raise HTTPException(503, "GEMINI_API_KEY не настроен на сервере")
        log.info("[GEMINI] отправка запроса...")
        import google.generativeai as genai
        genai.configure(api_key=GEMINI_API_KEY)
        m = genai.GenerativeModel("gemini-2.5-flash")
        result = m.generate_content(prompt).text

    elapsed = round(time.time() - t0, 1)
    log.info(f"[TRANSLATE] готово за {elapsed}с — символов в ответе={len(result)}")
    return result

# ── Extractors ───────────────────────────────────────────────────────────────
def extract_text_from_pdf(b):
    log.info("[EXTRACT] PDF")
    doc = fitz.open(stream=b, filetype="pdf")
    return "\n\n".join(p.get_text() for p in doc).strip()

def extract_text_from_docx(b):
    log.info("[EXTRACT] DOCX")
    doc = Document(io.BytesIO(b))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip()).strip()

def extract_text_from_pptx(b):
    log.info("[EXTRACT] PPTX")
    prs = Presentation(io.BytesIO(b))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
    return "\n".join(texts).strip()

def make_docx(text):
    doc = Document()
    for p in text.split("\n"):
        if p.strip():
            doc.add_paragraph(p.strip())
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()

# ── Routes ───────────────────────────────────────────────────────────────────
@app.get("/")
def root():
    return {"status": "ok"}

@app.get("/stats")
def get_stats():
    return load_stats()

@app.post("/visit")
def record_visit():
    return increment("visits")

class TextRequest(BaseModel):
    text: str
    model: Optional[str] = "gemini"

@app.post("/translate-text")
def translate_text(req: TextRequest):
    if not req.text.strip():
        raise HTTPException(422, "Текст пустой")
    log.info(f"[/translate-text] модель={req.model} символов={len(req.text)}")
    translated = translate_to_kazakh(req.text, req.model or "gemini")
    increment("translations")
    return {"translated": translated}

@app.post("/translate")
async def translate_document(
    file: UploadFile = File(...),
    model: Optional[str] = Form("gemini")
):
    name = file.filename.lower()
    log.info(f"[/translate] файл={file.filename} модель={model}")
    if not (name.endswith(".pdf") or name.endswith(".docx") or name.endswith(".txt") or name.endswith(".pptx")):
        raise HTTPException(400, "Только PDF, DOCX, PPTX, TXT")

    t0 = time.time()
    b = await file.read()
    log.info(f"[/translate] размер файла={round(len(b)/1024, 1)}КБ")

    if name.endswith(".pdf"):
        text = extract_text_from_pdf(b)
    elif name.endswith(".docx"):
        text = extract_text_from_docx(b)
    elif name.endswith(".pptx"):
        text = extract_text_from_pptx(b)
    else:
        text = b.decode("utf-8")

    log.info(f"[/translate] извлечено символов={len(text)}")
    if not text.strip():
        raise HTTPException(422, "Файл пустой")

    translated = translate_to_kazakh(text, model or "gemini")
    increment("translations")

    docx_bytes = make_docx(translated)
    elapsed = round(time.time() - t0, 1)
    log.info(f"[/translate] итого времени={elapsed}с")

    return StreamingResponse(
        io.BytesIO(docx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": "attachment; filename=translated_kazakh.docx"}
    )
